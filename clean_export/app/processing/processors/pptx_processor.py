import tempfile
from typing import AsyncGenerator
from app.processing.base import DocumentProcessor, ProcessingResult
from core.exceptions.processing import ProcessingFailedException

class PptxProcessor(DocumentProcessor):
    def supports(self, mime_type: str) -> bool:
        return mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    async def process(self, file_stream: AsyncGenerator[bytes, None]) -> ProcessingResult:
        with tempfile.SpooledTemporaryFile(max_size=5_000_000, mode='w+b') as spooled_file:
            async for chunk in file_stream:
                spooled_file.write(chunk)
                
            spooled_file.seek(0)
            
            try:
                import pptx
            except ImportError:
                raise ProcessingFailedException("python-pptx is required for PPTX extraction but is not installed.")
                
            try:
                presentation = pptx.Presentation(spooled_file)
                
                raw_text_parts = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            raw_text_parts.append(shape.text)
                            
                raw_text = "\n\n".join(raw_text_parts)
                
                character_count = len(raw_text)
                word_count = len(raw_text.split())
                slide_count = len(presentation.slides)
                
                document_metadata = {}
                if presentation.core_properties:
                    props = presentation.core_properties
                    document_metadata = {
                        "author": props.author,
                        "title": props.title,
                        "subject": props.subject,
                        "category": props.category,
                        "keywords": props.keywords,
                        "slide_count": slide_count
                    }
                    document_metadata = {k: v for k, v in document_metadata.items() if v}
                    
                processing_metadata = {
                    "parser_library": "python-pptx",
                    "parser_version": pptx.__version__
                }
                
                return ProcessingResult(
                    raw_text=raw_text,
                    page_count=slide_count,
                    word_count=word_count,
                    character_count=character_count,
                    document_metadata=document_metadata,
                    processing_metadata=processing_metadata,
                    detected_language=None
                )
            except Exception as e:
                raise ProcessingFailedException(f"Failed to process PPTX: {str(e)}")
